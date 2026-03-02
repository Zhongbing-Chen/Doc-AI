#!/usr/bin/env python3
"""
Test script: Parse PDF with VLM and output JSON

Usage:
    python test_vlm_parser.py [pdf_file] [vlm_url]

Example:
    python test_vlm_parser.py /path/to/document.pdf http://localhost:8000
"""

import sys
import json
from pathlib import Path

# Use package import
from document_parser import DocumentParser, Document


def test_with_vlm(pdf_path: str, vlm_url: str):
    """Parse PDF using VLM and print JSON result"""

    print("=" * 60)
    print("Document Parser - VLM Test")
    print("=" * 60)
    print(f"\nPDF File: {pdf_path}")
    print(f"VLM URL: {vlm_url}")
    print()

    # Create parser with VLM
    parser = DocumentParser(vlm_url=vlm_url)

    # Parse the PDF
    print("Parsing document with VLM...")
    doc = parser.parse(pdf_path)

    # Print metadata
    print("\n--- Document Metadata ---")
    print(f"File: {doc.metadata.file_name}")
    print(f"Type: {doc.metadata.file_type}")
    print(f"Size: {doc.metadata.file_size} bytes")
    print(f"Total Pages: {doc.document_info.total_pages}")
    print(f"Title: {doc.document_info.title or 'N/A'}")

    # Print page summaries
    print("\n--- Pages ---")
    for i, page in enumerate(doc.pages):
        print(f"\nPage {page.page_num}:")
        print(f"  Size: {page.width} x {page.height}")
        print(f"  Blocks: {len(page.blocks)}")
        for block in page.blocks[:3]:  # Show first 3 blocks
            print(f"    - {block.type}: {block.content.text[:50]}..." if len(block.content.text) > 50 else f"    - {block.type}: {block.content.text}")
        if len(page.blocks) > 3:
            print(f"    ... and {len(page.blocks) - 3} more blocks")

    # Output full JSON
    print("\n--- Full JSON Output ---")
    json_str = doc.to_json(indent=2)
    print(json_str)

    # Save to file
    output_path = Path("/tmp/parsed_document.json")
    doc.save(output_path)
    print(f"\n--- JSON saved to: {output_path} ---")

    return doc


def test_without_vlm(pdf_path: str):
    """Parse PDF using standard processor (no VLM)"""

    print("=" * 60)
    print("Document Parser - Standard Test (No VLM)")
    print("=" * 60)
    print(f"\nPDF File: {pdf_path}")
    print()

    # Create parser without VLM
    parser = DocumentParser()

    # Parse the PDF
    print("Parsing document...")
    doc = parser.parse(pdf_path)

    # Print metadata
    print("\n--- Document Metadata ---")
    print(f"File: {doc.metadata.file_name}")
    print(f"Type: {doc.metadata.file_type}")
    print(f"Size: {doc.metadata.file_size} bytes")
    print(f"Total Pages: {doc.document_info.total_pages}")

    # Print page summaries
    print("\n--- Pages ---")
    for i, page in enumerate(doc.pages):
        print(f"\nPage {page.page_num}:")
        print(f"  Blocks: {len(page.blocks)}")
        for block in page.blocks[:3]:
            text = block.content.text[:50] + "..." if len(block.content.text) > 50 else block.content.text
            print(f"    - {block.type}: {text}")
        if len(page.blocks) > 3:
            print(f"    ... and {len(page.blocks) - 3} more blocks")

    # Output JSON
    print("\n--- JSON Output (first 2000 chars) ---")
    json_str = doc.to_json(indent=2)
    print(json_str[:2000])
    if len(json_str) > 2000:
        print(f"\n... ({len(json_str) - 2000} more characters)")

    # Save to file
    output_path = Path("/tmp/parsed_document.json")
    doc.save(output_path)
    print(f"\n--- JSON saved to: {output_path} ---")

    return doc


def test_other_formats():
    """Test parsing other document formats"""

    print("=" * 60)
    print("Document Parser - Other Formats Test")
    print("=" * 60)

    # Create test files
    test_files = []

    # Word
    try:
        from docx import Document
        doc = Document()
        doc.add_heading('Test Word', 0)
        doc.add_paragraph('This is a test paragraph.')
        doc.save('/tmp/test.docx')
        test_files.append(('word', '/tmp/test.docx'))
    except Exception as e:
        print(f"Could not create Word test: {e}")

    # Excel
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws['A1'] = 'Name'
        ws['B1'] = 'Value'
        ws['A2'] = 'Item 1'
        ws['B2'] = 100
        wb.save('/tmp/test.xlsx')
        test_files.append(('excel', '/tmp/test.xlsx'))
    except Exception as e:
        print(f"Could not create Excel test: {e}")

    # PowerPoint
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test PPT"
        prs.save('/tmp/test.pptx')
        test_files.append(('powerpoint', '/tmp/test.pptx'))
    except Exception as e:
        print(f"Could not create PPT test: {e}")

    # Markdown
    Path('/tmp/test.md').write_text("# Test Markdown\n\nThis is test content.\n\n- Item 1\n- Item 2\n")
    test_files.append(('markdown', '/tmp/test.md'))

    # Parse each
    parser = DocumentParser()

    for file_type, file_path in test_files:
        print(f"\n--- Parsing {file_type.upper()} ---")
        doc = parser.parse(file_path)
        print(f"Pages: {doc.document_info.total_pages}")
        print(f"Blocks: {sum(len(p.blocks) for p in doc.pages)}")
        print(f"JSON size: {len(doc.to_json())} chars")

        # Save
        output = Path(f"/tmp/{file_type}_output.json")
        doc.save(output)
        print(f"Saved to: {output}")


def main():
    if len(sys.argv) > 2:
        # Parse PDF with VLM
        pdf_path = sys.argv[1]
        vlm_url = sys.argv[2]
        test_with_vlm(pdf_path, vlm_url)
    elif len(sys.argv) > 1:
        # Parse PDF without VLM
        pdf_path = sys.argv[1]
        if Path(pdf_path).exists():
            test_without_vlm(pdf_path)
        else:
            print(f"File not found: {pdf_path}")
            return 1
    else:
        # Test with sample files
        print("No arguments provided. Running format tests...")
        test_other_formats()

        # Try to find a PDF in the current directory
        for pdf in Path('.').glob('*.pdf'):
            print(f"\nFound PDF: {pdf}")
            test_without_vlm(str(pdf))
            break

    return 0


if __name__ == '__main__':
    sys.exit(main())
